#!/usr/bin/env python3
"""Build a detailed product/function report and PPT for Hermass Observer."""

from __future__ import annotations

import argparse
import html
import json
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "outputs" / "system_function_report"
PUBLIC_DIR = ROOT / "public"


def ymd(date_str: str) -> str:
    return date_str.replace("-", "")


def load_json(path: Path) -> dict[str, Any]:
    return json.loads(path.read_text(encoding="utf-8"))


def pct(value: float | None, digits: int = 1) -> str:
    if value is None:
        return "-"
    return f"{value * 100:.{digits}f}%"


def pnum(value: Any) -> str:
    if value is None:
        return "-"
    if isinstance(value, float):
        return f"{value:,.2f}"
    if isinstance(value, int):
        return f"{value:,}"
    return str(value)


def collect_context(date_str: str) -> dict[str, Any]:
    date_ymd = ymd(date_str)
    paths = {
        "brief": ROOT / "outputs" / "daily_research_brief" / f"daily_research_brief_{date_ymd}.json",
        "state_manifest": ROOT / "outputs" / "state_cache" / f"state_cache_manifest_{date_ymd}.json",
        "signals": ROOT / "outputs" / "strategy_signals" / f"strategy_signal_daily_{date_ymd}.json",
        "fit_log": ROOT / "outputs" / "strategy_fit_observer" / f"fit_log_{date_ymd}.json",
        "forward": ROOT / "outputs" / "forward_observation" / f"forward_observation_{date_ymd}.json",
        "rotation": ROOT / "outputs" / "industry_rotation" / f"industry_rotation_{date_ymd}.json",
        "ifind": ROOT / "outputs" / "ifind" / f"financial_{date_ymd}.json",
    }
    data = {key: load_json(path) for key, path in paths.items()}
    data["paths"] = {key: str(path) for key, path in paths.items()}
    return data


def build_markdown(date_str: str, ctx: dict[str, Any]) -> str:
    brief = ctx["brief"]
    state = ctx["state_manifest"]
    signals = ctx["signals"]
    fit = ctx["fit_log"]
    forward = ctx["forward"]
    rotation = ctx["rotation"]
    ifind = ctx["ifind"]
    market = brief["market_summary"]
    stats = brief["signal_stats"]
    top_industries = rotation.get("top_industries", [])[:8]

    lines = [
        "# Hermass Observer 功能说明与投研总报系统白皮书",
        "",
        f"**报告日期：{date_str}**",
        f"**生成时间：{datetime.now(timezone.utc).isoformat()}**",
        "",
        "## 1. 一句话定位",
        "",
        "Hermass Observer 是一个面向 A 股的多周期 State 观察、策略信号账本、基本面场景过滤和投研总报系统。它不试图预测每只股票的涨跌，而是回答一个更可执行的问题：",
        "",
        "> 今天，哪些股票在适合它们策略的环境里触发了信号？",
        "",
        "系统核心链条是：",
        "",
        "```text",
        "State 底座寻找先验概率更大的环境",
        "  -> 经典策略模块提供精确触发",
        "  -> 策略信号账本记录事实",
        "  -> 场景过滤层标注生命周期和适配度",
        "  -> iFinD 提供公司质量和产业身份背景",
        "  -> 投研总报只呈现事实与背景，不输出交易指令",
        "```",
        "",
        "## 2. 设计哲学",
        "",
        "### 2.1 模糊的正确",
        "",
        "系统不追求用复杂模型给出精确预测，而是通过多周期共振、收缩/扩张、波动状态、行业集中度和公司质量，把不适合趋势策略的场景先过滤掉。",
        "",
        "### 2.2 底座不污染",
        "",
        "State 底座只负责计算市场状态，不被策略、推荐、基本面或 LLM 逻辑污染。所有增强都发生在缓存层、账本层、观察层和报告层。",
        "",
        "### 2.3 策略不重写",
        "",
        "VCP、2560、布林强盗各自独立输出权威信号。提醒层和总报不重新计算策略触发条件，只消费 `strategy_signal_daily` 中的标准化事实。",
        "",
        "### 2.4 诚实反馈",
        "",
        "校准没通过就显示“待校准”。前向观察没有未来行情就显示 `pending_future_data`。iFinD 无明确字段就留空或标注“未标注”。",
        "",
        "## 3. 系统分层架构",
        "",
        "| 层级 | 模块 | 职责 | 主要产物 |",
        "|------|------|------|----------|",
        "| L0 | 黑狼日线/资金流/iFinD 数据 | 原始事实采集 | 日线 zip/raw DB、资金流 DB、fundamental_evidence.duckdb |",
        "| L1 | Foundation + State 底座 | 计算 MN1/W1/D1 state、SR、指标 | p116_foundation.duckdb |",
        "| L2 | State Cache | 把慢查询缓存成小 JSON/表 | state_ef、state_duration、sr_boundary |",
        "| L3 | 策略信号账本 | 标准化 VCP/2560/布林强盗信号 | strategy_signal_daily |",
        "| L4 | 场景过滤与观察 | 生命周期、适配度、前向样本 | fit_log、forward_observation |",
        "| L5 | 投研总报 | 面向订阅者呈现事实 | daily_research_brief.html/md/json |",
        "",
        "## 4. 每日 Agently 链路",
        "",
        "2026-05-22 已验证的每日链路包括：",
        "",
        "1. `preflight_freshness`：检查脚本和工作流新鲜度。",
        "2. `download_daily`：下载并合并黑狼日线。",
        "3. `build_raw_db`：构建 P108 raw DuckDB。",
        "4. `download_moneyflow` + `import_moneyflow_db`：下载并导入资金流。",
        "5. `build_moneyflow_evidence`：生成 5 日资金流证据。",
        "6. `download_market_assets` + `build_market_assets_state`：更新指数/ETF 状态。",
        "7. `build_foundation`：构建 P116 foundation。",
        "8. `build_state_cache`：生成 State 快照、持续天数、SR 边界。",
        "9. `build_strategy_signal_ledger`：生成策略信号账本。",
        "10. `evaluate_strategy_evidence`：生成策略证据分解。",
        "11. `build_strategy_reminder`：生成策略提醒。",
        "12. `process_ifind_data`：标准化 iFinD 公司质量和产业身份。",
        "13. `run_strategy_fit_observer`：保存策略-环境适配观察。",
        "14. `forward_sim`：生成前向观察样本。",
        "15. `generate_daily_brief`：生成每日投研总报。",
        "",
        "## 5. 2026-05-22 真实运行结果",
        "",
        "| 指标 | 数值 |",
        "|------|------|",
        f"| 全三 E/F 池 | {market['all_three_ef_count']} 只 |",
        f"| 较上一缓存日变化 | {market['all_three_ef_delta']:+d} |",
        f"| State cache: duration rows | {state['counts']['duration_rows']} |",
        f"| State cache: SR boundary rows | {state['counts']['sr_boundary_rows']} |",
        f"| 策略信号总数 | {signals['signal_count']} |",
        f"| 正式提醒信号 | {stats['total_reminders']} |",
        f"| 总报主体展示 | {len(brief['display_rows'])} |",
        f"| iFinD 场景聚焦 | {len(brief['ifind_focus_rows'])} |",
        f"| 前向观察样本 | {forward['total']} |",
        f"| 前向观察已标注 | {forward['labeled']} |",
        f"| 前向观察待更新 | {forward['pending']} |",
        "",
        "### 5.1 策略信号分布",
        "",
        "| 策略信号 | 数量 |",
        "|----------|------|",
    ]
    for key, value in signals["strategy_counts"].items():
        lines.append(f"| {key} | {value} |")

    lines += [
        "",
        "### 5.2 适配度分布",
        "",
        "| 适配度 | 数量 |",
        "|--------|------|",
    ]
    for key, value in fit["fit_counts"].items():
        lines.append(f"| {key} | {value} |")

    lines += [
        "",
        "### 5.3 生命周期分布",
        "",
        "| 生命周期 | 数量 |",
        "|----------|------|",
    ]
    for key, value in fit["lifecycle_counts"].items():
        lines.append(f"| {key} | {value} |")

    lines += [
        "",
        "### 5.4 行业轮动 Top 8",
        "",
        "| 排名 | 行业 | E/F 池数量 | 新进 | 资金确认率 | 轮动分 |",
        "|------|------|-------------|------|------------|--------|",
    ]
    for row in top_industries:
        lines.append(
            f"| {row['rank']} | {row['sw_l1']} | {row['pool_count']} | {row['entered_count']} | {pct(row.get('moneyflow_confirm_rate'))} | {pnum(row.get('rotation_score'))} |"
        )

    lines += [
        "",
        "### 5.5 iFinD 公司质量覆盖",
        "",
        f"- 覆盖股票：{ifind['total']} 只",
        f"- 质量健康：{ifind['quality_counts'].get('质量健康', 0)}",
        f"- 质量中性：{ifind['quality_counts'].get('质量中性', 0)}",
        f"- 质量谨慎：{ifind['quality_counts'].get('质量谨慎', 0)}",
        "",
        "## 6. 核心功能详解",
        "",
        "### 6.1 State 底座",
        "",
        "State 底座是系统的核心先验环境识别器。它从 OHLCV、趋势、位置、波动等状态编码出 MN1/W1/D1 三周期状态，形成全三 E/F 共振池。",
        "",
        "关键能力：",
        "",
        "- 识别多周期趋势共振。",
        "- 标注 D1/W1/MN1 是否刚脱离收缩。",
        "- 记录三周期共振持续天数。",
        "- 生成 SR 边界和价格相对边界位置。",
        "- 把全市场慢查询缓存成毫秒级 JSON。",
        "",
        "### 6.2 策略信号账本",
        "",
        "策略信号账本是策略模块与提醒层之间的契约。它记录每个策略模块的正式输出，并统一字段。",
        "",
        "关键字段：",
        "",
        "- `strategy_id`: vcp / ma2560 / bollinger_bandit",
        "- `signal_type`: entry / exit / risk / structure",
        "- `reminder_eligible`: 是否可进入提醒层",
        "- `display_scope`: reminder / research",
        "- `lifecycle_stage`: 新生 / 行进 / 延展 / 未知",
        "- `strategy_environment_fit`: 最佳适配 / 适配 / 弱适配 / 待观察",
        "- `fit_reasons`: 适配依据",
        "",
        "### 6.3 场景过滤",
        "",
        "场景过滤不是策略，不产生信号。它把“这个信号出现时的背景”结构化，方便未来按“策略 × 生命周期 × 适配度”做统计。",
        "",
        "示例：",
        "",
        "- VCP + 新生 + 三周期共振新近形成 -> 最佳适配",
        "- 2560 + 行进 + 波动稳定 -> 最佳适配",
        "- 布林强盗 + 延展 + D1波动偏活跃 -> 最佳适配",
        "",
        "### 6.4 iFinD 基本面场景",
        "",
        "iFinD 数据只做背景标注和场景过滤，不参与策略触发。",
        "",
        "当前接入字段：",
        "",
        "- 公司质量分",
        "- 主业纯度",
        "- 现金含金量",
        "- 盈利质量",
        "- 资产安全",
        "- 申万一级/二级/三级行业",
        "- 同花顺概念",
        "- 主营业务和主要产品",
        "",
        "### 6.5 每日投研总报",
        "",
        "总报回答订阅者最关心的问题：今天哪些股票在适合它们策略的环境里触发了信号？",
        "",
        "报告主体只展示：",
        "",
        "- 最佳适配",
        "- 适配",
        "",
        "弱适配和待观察进入统计摘要，不进入主体表，降低噪音。",
        "",
        "### 6.6 前向观察",
        "",
        "前向观察账本不是交易模拟器。它记录正式提醒信号在未来 5/10/20 个交易日的表现，等未来行情可用后自动补标签。",
        "",
        "它不模拟成交、不推断离场、不做仓位。",
        "",
        "## 7. 主要输出文件",
        "",
        "| 类型 | 完整路径 |",
        "|------|----------|",
        f"| 投研总报 HTML | `{ROOT / 'public' / f'daily_research_brief_{ymd(date_str)}.html'}` |",
        f"| 投研总报 Markdown | `{ROOT / 'outputs' / 'daily_research_brief' / f'daily_research_brief_{ymd(date_str)}.md'}` |",
        f"| 策略提醒 HTML | `{ROOT / 'public' / f'strategy_reminder_{ymd(date_str)}.html'}` |",
        f"| 前向观察 HTML | `{ROOT / 'public' / f'forward_observation_{ymd(date_str)}.html'}` |",
        f"| 行业轮动 HTML | `{ROOT / 'public' / f'industry_rotation_{ymd(date_str)}.html'}` |",
        f"| 全三 E/F 池 HTML | `{ROOT / 'public' / f'p116_all_three_ef_{ymd(date_str)}.html'}` |",
        "",
        "## 8. 合规边界",
        "",
        "系统所有对外报告遵守以下边界：",
        "",
        "- 不输出具体交易指令。",
        "- 不使用未通过校准的胜率数字。",
        "- 不把 iFinD 基本面数据用于策略触发。",
        "- 不用 LLM 编造事实。",
        "- 缺数据时明确留空或标注待校准/待更新。",
        "",
        "## 9. 后续路线图",
        "",
        "1. 累积前向观察样本，补齐未来 5/10/20 日表现。",
        "2. 按“策略 × 生命周期 × 适配度 × iFinD 质量标签”做统计。",
        "3. 将每日总报推送到固定渠道。",
        "4. 引入人工交易日志，对比系统信号与实际执行损耗。",
        "5. 等校准质量闸门通过后，在报告中展示真实历史统计。",
        "",
    ]
    return "\n".join(lines)


def build_html(markdown: str, date_str: str) -> str:
    body = html.escape(markdown)
    # Minimal Markdown-ish rendering for local viewing.
    lines = []
    in_code = False
    for raw in markdown.splitlines():
        line = html.escape(raw)
        if raw.startswith("```"):
            lines.append("</pre>" if in_code else "<pre>")
            in_code = not in_code
        elif in_code:
            lines.append(line)
        elif raw.startswith("# "):
            lines.append(f"<h1>{line[2:]}</h1>")
        elif raw.startswith("## "):
            lines.append(f"<h2>{line[3:]}</h2>")
        elif raw.startswith("### "):
            lines.append(f"<h3>{line[4:]}</h3>")
        elif raw.startswith("- "):
            lines.append(f"<p class='bullet'>• {line[2:]}</p>")
        elif raw.startswith("|"):
            lines.append(f"<p class='tableline'>{line}</p>")
        elif raw.strip() == "":
            lines.append("")
        else:
            lines.append(f"<p>{line}</p>")
    return f"""<!doctype html>
<html lang="zh-CN">
<head>
  <meta charset="utf-8">
  <meta name="viewport" content="width=device-width, initial-scale=1">
  <title>Hermass Observer 功能说明 {date_str}</title>
  <style>
    body {{ margin: 0; background: #f5f7fb; color: #172033; font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", sans-serif; }}
    main {{ max-width: 1080px; margin: 0 auto; padding: 34px 28px 80px; background: white; min-height: 100vh; }}
    h1 {{ font-size: 32px; margin: 0 0 18px; color: #101828; }}
    h2 {{ font-size: 22px; margin: 34px 0 12px; color: #1d2939; border-top: 1px solid #eaecf0; padding-top: 22px; }}
    h3 {{ font-size: 17px; margin: 22px 0 8px; color: #344054; }}
    p {{ font-size: 14px; line-height: 1.72; margin: 6px 0; }}
    .bullet {{ padding-left: 18px; }}
    .tableline {{ font-family: ui-monospace, SFMono-Regular, Menlo, monospace; white-space: pre-wrap; background: #f8fafc; padding: 3px 8px; border-radius: 4px; overflow-x: auto; }}
    pre {{ background: #101828; color: #e4e7ec; padding: 14px; border-radius: 8px; overflow-x: auto; }}
  </style>
</head>
<body><main>{chr(10).join(lines)}</main></body>
</html>"""


def add_textbox(slide, x, y, w, h, text, size=18, bold=False, color=RGBColor(23, 32, 51)):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    return box


def add_bullets(slide, x, y, w, h, bullets, size=15):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    for idx, item in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = RGBColor(52, 64, 84)
        p.space_after = Pt(6)
    return box


def add_table(slide, x, y, w, h, rows):
    table_shape = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(val)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(10 if r else 10)
                paragraph.font.bold = r == 0
                paragraph.font.color.rgb = RGBColor(23, 32, 51)
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(234, 240, 248)
    return table_shape


def build_pptx(date_str: str, ctx: dict[str, Any], out_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def title_slide(title, subtitle):
        slide = prs.slides.add_slide(blank)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(245, 247, 251)
        add_textbox(slide, 0.75, 1.0, 11.8, 1.0, title, size=34, bold=True, color=RGBColor(16, 24, 40))
        add_textbox(slide, 0.78, 2.0, 11.5, 0.7, subtitle, size=18, color=RGBColor(71, 84, 103))
        add_textbox(
            slide,
            0.78,
            6.8,
            10,
            0.3,
            f"Report date: {date_str} | Research-only system overview",
            size=10,
            color=RGBColor(102, 112, 133),
        )

    def section(title, bullets):
        slide = prs.slides.add_slide(blank)
        add_textbox(slide, 0.55, 0.35, 12.2, 0.45, title, size=24, bold=True)
        add_bullets(slide, 0.75, 1.15, 11.8, 5.8, bullets, size=16)

    brief = ctx["brief"]
    state = ctx["state_manifest"]
    signals = ctx["signals"]
    fit = ctx["fit_log"]
    forward = ctx["forward"]
    rotation = ctx["rotation"]
    ifind = ctx["ifind"]
    market = brief["market_summary"]
    stats = brief["signal_stats"]

    title_slide(
        "Hermass Observer",
        "多周期 State 观察、策略信号账本、iFinD 场景过滤与每日投研总报系统",
    )

    section(
        "系统回答的问题",
        [
            "不是预测每只股票涨跌，而是筛出适合趋势策略的高先验环境。",
            "核心问题：今天哪些股票在适合它们策略的环境里触发了信号？",
            "State 负责环境，策略负责触发，iFinD 负责公司/产业背景，总报负责呈现事实。",
            "校准未通过时保持“待校准”，不展示不可靠胜率。",
        ],
    )

    section(
        "核心哲学",
        [
            "模糊的正确：先排除不适合趋势交易的场景。",
            "底座不污染：State 只做确定性计算，不混入策略和 LLM。",
            "策略不重写：提醒层只消费 VCP/2560/布林强盗模块的标准化输出。",
            "诚实反馈：缺数据、未校准、未完成前向标签时明确显示状态。",
        ],
    )

    section(
        "分层架构",
        [
            "L0 数据采集：黑狼日线、资金流、市场资产、iFinD。",
            "L1 Foundation：OHLCV、SR、MN1/W1/D1 State。",
            "L2 State Cache：全三 E/F、持续天数、SR 边界、转换。",
            "L3 策略信号账本：VCP、2560、布林强盗的权威信号。",
            "L4 场景过滤：生命周期、环境适配、适配理由、前向观察。",
            "L5 投研总报：策略-股票-环境-iFinD 背景的一页式呈现。",
        ],
    )

    slide = prs.slides.add_slide(blank)
    add_textbox(slide, 0.55, 0.35, 12, 0.45, "2026-05-22 真实运行数字", size=24, bold=True)
    rows = [
        ["指标", "数值"],
        ["全三 E/F 池", market["all_three_ef_count"]],
        ["较上一缓存日", f"{market['all_three_ef_delta']:+d}"],
        ["策略信号总数", signals["signal_count"]],
        ["正式提醒信号", stats["total_reminders"]],
        ["总报主体展示", len(brief["display_rows"])],
        ["iFinD 场景聚焦", len(brief["ifind_focus_rows"])],
        ["前向观察样本", forward["total"]],
    ]
    add_table(slide, 0.8, 1.2, 5.4, 4.8, rows)
    add_bullets(
        slide,
        6.7,
        1.2,
        5.8,
        4.5,
        [
            f"State cache 持续天数行数：{state['counts']['duration_rows']}",
            f"SR 边界行数：{state['counts']['sr_boundary_rows']}",
            f"前向观察待更新：{forward['pending']}",
            f"校准状态：{brief['calibration']['status']} / {brief['calibration']['reason']}",
        ],
        size=15,
    )

    slide = prs.slides.add_slide(blank)
    add_textbox(slide, 0.55, 0.35, 12, 0.45, "策略信号与场景适配", size=24, bold=True)
    rows = [["类别", "数量"]]
    rows.extend([[k, v] for k, v in signals["strategy_counts"].items()])
    add_table(slide, 0.65, 1.15, 5.6, 4.9, rows)
    rows2 = [["适配度", "数量"]]
    rows2.extend([[k, v] for k, v in fit["fit_counts"].items()])
    add_table(slide, 6.7, 1.15, 5.6, 3.1, rows2)

    section(
        "三套策略如何与 State 匹配",
        [
            "VCP：趋势新生，D1 刚脱离收缩或三周期共振新近形成。",
            "2560：趋势行进，D1 E/F 持续、波动稳定、趋势结构稳定。",
            "布林强盗：趋势延展，波动从稳定转活跃或价格位于阻力区间上方。",
            "适配字段不是交易建议，只是信号出现时的环境解释。",
        ],
    )

    slide = prs.slides.add_slide(blank)
    add_textbox(slide, 0.55, 0.35, 12, 0.45, "行业轮动与资金流", size=24, bold=True)
    rows = [["排名", "行业", "E/F池", "新进", "资金确认率", "轮动分"]]
    for row in rotation.get("top_industries", [])[:7]:
        rows.append(
            [
                row["rank"],
                row["sw_l1"],
                row["pool_count"],
                row["entered_count"],
                pct(row.get("moneyflow_confirm_rate")),
                pnum(row.get("rotation_score")),
            ]
        )
    add_table(slide, 0.55, 1.1, 12.1, 5.6, rows)

    section(
        "iFinD 基本面与产业身份",
        [
            f"覆盖股票：{ifind['total']} 只。",
            f"质量健康：{ifind['quality_counts'].get('质量健康', 0)}；质量中性：{ifind['quality_counts'].get('质量中性', 0)}；质量谨慎：{ifind['quality_counts'].get('质量谨慎', 0)}。",
            "接入字段包括：质量分、现金含金量、盈利质量、资产安全、申万行业、概念、主营业务、主要产品。",
            "iFinD 只做公司质地和产业背景标注，不参与策略触发。",
        ],
    )

    section(
        "每日投研总报如何使用",
        [
            "打开 daily_research_brief_latest.html。",
            "先看市场环境速览：全三 E/F 池规模、行业聚焦、校准状态。",
            "再看最佳适配与适配信号：股票、策略、生命周期、适配理由、State、SR、iFinD 背景。",
            "弱适配和待观察不进入主体表，只在统计中保留。",
        ],
    )

    section(
        "合规边界",
        [
            "报告只呈现策略信号与环境匹配事实。",
            "不输出具体交易指令。",
            "不展示未通过校准的胜率。",
            "不使用 LLM 编造公司、行业或统计数据。",
            "所有缺失项保持空值、未标注、待校准或待更新。",
        ],
    )

    section(
        "后续路线图",
        [
            "持续积累前向观察样本，补齐未来 5/10/20 日标签。",
            "按策略 × 生命周期 × 适配度 × iFinD 质量标签做统计。",
            "接入人工交易日志，评估真实执行损耗。",
            "质量闸门通过后，在总报中展示真实历史参考。",
            "将总报推送到固定渠道，形成每日 5 分钟投研流程。",
        ],
    )

    prs.save(out_path)


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--date", required=True)
    args = parser.parse_args()
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    PUBLIC_DIR.mkdir(parents=True, exist_ok=True)
    ctx = collect_context(args.date)
    date_ymd = ymd(args.date)
    markdown = build_markdown(args.date, ctx)
    md_path = OUT_DIR / f"hermass_observer_function_report_{date_ymd}.md"
    html_path = PUBLIC_DIR / f"hermass_observer_function_report_{date_ymd}.html"
    pptx_path = OUT_DIR / f"hermass_observer_function_report_{date_ymd}.pptx"
    md_path.write_text(markdown, encoding="utf-8")
    html_path.write_text(build_html(markdown, args.date), encoding="utf-8")
    build_pptx(args.date, ctx, pptx_path)
    latest_md = OUT_DIR / "hermass_observer_function_report_latest.md"
    latest_html = PUBLIC_DIR / "hermass_observer_function_report_latest.html"
    latest_pptx = OUT_DIR / "hermass_observer_function_report_latest.pptx"
    latest_md.write_text(markdown, encoding="utf-8")
    latest_html.write_text(build_html(markdown, args.date), encoding="utf-8")
    latest_pptx.write_bytes(pptx_path.read_bytes())
    print(
        json.dumps(
            {
                "ok": True,
                "date": args.date,
                "markdown": str(md_path),
                "html": str(html_path),
                "pptx": str(pptx_path),
                "latest_markdown": str(latest_md),
                "latest_html": str(latest_html),
                "latest_pptx": str(latest_pptx),
            },
            ensure_ascii=False,
            indent=2,
        )
    )
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
